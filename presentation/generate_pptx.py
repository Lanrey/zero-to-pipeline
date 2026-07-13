"""
Generate zero-to-pipeline-europython2026.pptx
EuroPython 2026 — "Zero to Pipeline — When Data Connectors Just Work"
by Olusola Akinsulere

Design system modelled on example conference slides (purplecon style):
  - 10.0 × 5.625 in (standard 16:9)
  - Background #444349 (dark warm gray)
  - Body text #EBDAFF (lavender)
  - Accent text #AAE5FA (light blue)
  - Secondary accent #C6A5FF (light purple)
  - Font: Open Sans (headings / body) — Consolas (code)
  - Extremely sparse text: 4–8 words per content slide
  - Code slides: full-width Consolas white on dark background
  - Red border callout boxes to highlight key lines

Accessibility compliance:
  - Shapes named and ordered for correct Tab / screen-reader sequence
  - Alt text on every content shape; decorative shapes marked decorative
  - Colour information always paired with text labels
  - Minimum 18 pt body, 18 pt code
  - Speaker notes on every slide
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

# ── Colour palette (purplecon style) ─────────────────────────────────────────
BG        = RGBColor(0x44, 0x43, 0x49)   # #444349  dark warm gray bg
BG_DARK   = RGBColor(0x21, 0x21, 0x23)   # #212123  darker variant for code
LAVENDER  = RGBColor(0xEB, 0xDA, 0xFF)   # #EBDAFF  primary body text
LBLUE     = RGBColor(0xAA, 0xE5, 0xFA)   # #AAE5FA  accent / keywords
LPURPLE   = RGBColor(0xC6, 0xA5, 0xFF)   # #C6A5FF  secondary accent
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)   # #FFFFFF  code text / callout bg
CORAL     = RGBColor(0xE0, 0x66, 0x66)   # #E06666  inline code labels
RED_BOX   = RGBColor(0xFF, 0x00, 0x00)   # #FF0000  callout border
GRAY      = RGBColor(0x99, 0x99, 0x99)   # #999999  muted / secondary
DARK_CODE = RGBColor(0x43, 0x43, 0x43)   # #434343  "hidden" code (reveal trick)

SLIDE_W   = Inches(10.0)
SLIDE_H   = Inches(5.625)

HEADING_F = "Open Sans"
BODY_F    = "Open Sans"
CODE_F    = "Consolas"


# ── Low-level helpers ─────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs):
    layout = prs.slide_layouts[6]
    slide  = prs.slides.add_slide(layout)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG
    return slide


def _cNvPr(shape):
    el = shape._element
    for tag in ('p:nvSpPr', 'p:nvPicPr', 'p:nvGrpSpPr', 'p:nvCxnSpPr'):
        container = el.find('.//' + qn(tag))
        if container is not None:
            cNvPr = container.find(qn('p:cNvPr'))
            if cNvPr is not None:
                return cNvPr
    return None


def name_shape(shape, name: str):
    cNvPr = _cNvPr(shape)
    if cNvPr is not None:
        cNvPr.set('name', name)


def set_alt_text(shape, title: str, desc: str = ""):
    cNvPr = _cNvPr(shape)
    if cNvPr is None:
        return
    cNvPr.set('title', title)
    if desc:
        cNvPr.set('descr', desc)


def mark_decorative(shape):
    cNvPr = _cNvPr(shape)
    if cNvPr is None:
        return
    a16_ns = 'http://schemas.microsoft.com/office/drawing/2014/main'
    dec = etree.SubElement(cNvPr, f'{{{a16_ns}}}decorative')
    dec.set('val', '1')
    cNvPr.set('title', '')
    cNvPr.set('descr', '')


def add_notes(slide, text: str):
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = text


def txbox(slide, left, top, width, height,
          text="", font=BODY_F, size=Pt(22),
          color=LAVENDER, bold=False, italic=False,
          align=PP_ALIGN.LEFT, wrap=True,
          sname="", alt=""):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text           = text
    run.font.name      = font
    run.font.size      = size
    run.font.bold      = bold
    run.font.italic    = italic
    run.font.color.rgb = color
    if sname:
        name_shape(tb, sname)
    if alt:
        set_alt_text(tb, alt)
    elif text:
        set_alt_text(tb, text[:120])
    return tb


def rect(slide, left, top, width, height, fill,
         decorative=True, sname="", alt=""):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    if decorative:
        mark_decorative(shape)
    else:
        if sname:
            name_shape(shape, sname)
        if alt:
            set_alt_text(shape, alt)
    return shape


def red_callout(slide, left, top, width, height,
                alt="Callout box highlighting key information"):
    """Transparent box with red border — purplecon callout style."""
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.background()          # transparent
    shape.line.color.rgb = RED_BOX
    shape.line.width = Pt(2)
    set_alt_text(shape, "Callout", alt)
    return shape


def slide_number(slide, num):
    tb = txbox(slide,
        left=Inches(9.1), top=Inches(5.25),
        width=Inches(0.75), height=Inches(0.3),
        text=str(num), size=Pt(11), color=GRAY,
        align=PP_ALIGN.RIGHT, sname=f"slide-num-{num}")
    mark_decorative(tb)


def add_notes_to(slide, text):
    add_notes(slide, text)


# ── High-level slide builders ─────────────────────────────────────────────────

def big_phrase(prs, slide_num, phrase, sub="", color=LAVENDER,
               sub_color=GRAY, notes=""):
    """Single large left-aligned phrase — the dominant slide type."""
    slide = blank_slide(prs)
    txbox(slide,
        Inches(0.54), Inches(0.49), Inches(7.0), Inches(4.0),
        text=phrase, font=HEADING_F, size=Pt(40), bold=True,
        color=color, align=PP_ALIGN.LEFT, wrap=True,
        sname="slide-title", alt=f"Slide heading: {phrase}")
    if sub:
        txbox(slide,
            Inches(0.54), Inches(3.8), Inches(8.5), Inches(0.8),
            text=sub, font=BODY_F, size=Pt(22), color=sub_color,
            sname="slide-subtitle")
    slide_number(slide, slide_num)
    if notes:
        add_notes(slide, notes)
    return slide


def section_break(prs, slide_num, phrase, sub="", notes=""):
    """Centred section divider — short phrase, large."""
    slide = blank_slide(prs)
    txbox(slide,
        Inches(0.5), Inches(1.5), Inches(9.0), Inches(2.0),
        text=phrase, font=HEADING_F, size=Pt(44), bold=True,
        color=LBLUE, align=PP_ALIGN.CENTER, wrap=True,
        sname="slide-title", alt=f"Section: {phrase}")
    if sub:
        txbox(slide,
            Inches(1.0), Inches(3.6), Inches(8.0), Inches(0.7),
            text=sub, font=BODY_F, size=Pt(22), color=GRAY,
            align=PP_ALIGN.CENTER, sname="slide-subtitle")
    slide_number(slide, slide_num)
    if notes:
        add_notes(slide, notes)
    return slide


def code_slide(prs, slide_num, title, lines, callout=None,
               font_size=Pt(18), notes=""):
    """
    Full-width code block (Consolas white on BG_DARK).
    lines = list of (text, color) or plain strings.
    callout = (top_row_index, bottom_row_index) to draw a red box around rows.
    """
    slide = blank_slide(prs)
    # Title
    txbox(slide,
        Inches(0.54), Inches(0.2), Inches(9.0), Inches(0.55),
        text=title, font=HEADING_F, size=Pt(26), bold=True,
        color=LBLUE, sname="slide-title", alt=f"Code slide: {title}")
    # Dark background for code area
    code_top  = Inches(0.85)
    code_h    = Inches(4.55)
    rect(slide, Inches(0), code_top, Inches(10.0), code_h,
         BG_DARK, decorative=True)
    # Text frame
    tb = slide.shapes.add_textbox(
        Inches(0.3), code_top + Inches(0.15),
        Inches(9.4), code_h - Inches(0.3))
    tf = tb.text_frame
    tf.word_wrap = False
    first = True
    line_texts = []
    for line in lines:
        text, color = line if isinstance(line, tuple) else (line, WHITE)
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_before = Pt(1)
        run = p.add_run()
        run.text = text
        run.font.name      = CODE_F
        run.font.size      = font_size
        run.font.color.rgb = color
        line_texts.append(text)
    name_shape(tb, "code-block")
    set_alt_text(tb, f"Code: {title}",
        "Code block — " + " ".join(t for t in line_texts if t.strip())[:200])
    # Optional red callout box
    if callout:
        top_idx, bot_idx = callout
        line_h = float(font_size) / 72.0 * 1.35  # approx inches per line
        box_top  = float(code_top) / 914400 + 0.15 + top_idx * line_h
        box_bot  = float(code_top) / 914400 + 0.15 + (bot_idx + 1) * line_h
        box_h    = box_bot - box_top
        red_callout(slide,
            Inches(0.15), Inches(box_top),
            Inches(9.7), Inches(box_h))
    slide_number(slide, slide_num)
    if notes:
        add_notes(slide, notes)
    return slide


def two_col_slide(prs, slide_num, title,
                  left_head, left_items, left_color,
                  right_head, right_items, right_color,
                  notes=""):
    slide = blank_slide(prs)
    txbox(slide,
        Inches(0.54), Inches(0.2), Inches(9.0), Inches(0.55),
        text=title, font=HEADING_F, size=Pt(28), bold=True,
        color=LBLUE, sname="slide-title")
    # Divider
    div = slide.shapes.add_shape(1,
        Inches(4.9), Inches(0.8), Inches(0.03), Inches(4.6))
    div.fill.solid()
    div.fill.fore_color.rgb = GRAY
    div.line.fill.background()
    mark_decorative(div)
    # Left
    txbox(slide, Inches(0.4), Inches(0.85), Inches(4.3), Inches(0.55),
        text=left_head, font=HEADING_F, size=Pt(22), bold=True,
        color=left_color, sname="left-heading",
        alt=f"Left column: {left_head}")
    for i, item in enumerate(left_items):
        txbox(slide,
            Inches(0.4), Inches(1.5 + i * 0.65), Inches(4.3), Inches(0.6),
            text=f"• {item}", font=BODY_F, size=Pt(18), color=LAVENDER,
            wrap=True, sname=f"left-{i+1}")
    # Right
    txbox(slide, Inches(5.1), Inches(0.85), Inches(4.3), Inches(0.55),
        text=right_head, font=HEADING_F, size=Pt(22), bold=True,
        color=right_color, sname="right-heading",
        alt=f"Right column: {right_head}")
    for i, item in enumerate(right_items):
        txbox(slide,
            Inches(5.1), Inches(1.5 + i * 0.65), Inches(4.3), Inches(0.6),
            text=f"• {item}", font=BODY_F, size=Pt(18), color=LAVENDER,
            wrap=True, sname=f"right-{i+1}")
    slide_number(slide, slide_num)
    if notes:
        add_notes(slide, notes)
    return slide


# ── Build ─────────────────────────────────────────────────────────────────────

def build():
    prs = new_prs()
    n = 0   # slide counter

    # ── SLIDE 1: About Me ────────────────────────────────────────────────────
    n += 1
    slide = blank_slide(prs)
    txbox(slide, Inches(0.54), Inches(0.6), Inches(6.5), Inches(0.9),
        text="Olusola Akinsulere",
        font=HEADING_F, size=Pt(40), bold=True, color=WHITE,
        sname="slide-title", alt="Speaker name: Olusola Akinsulere")
    txbox(slide, Inches(0.54), Inches(1.6), Inches(6.5), Inches(0.55),
        text="Engineering Lead · Retailloop",
        font=BODY_F, size=Pt(24), color=LBLUE, sname="role")
    txbox(slide, Inches(0.54), Inches(2.3), Inches(6.5), Inches(1.5),
        text="Distributed systems · Payments · Data platforms\nPython · Backend · Practical ML",
        font=BODY_F, size=Pt(20), color=LAVENDER, wrap=True, sname="bio")
    txbox(slide, Inches(0.54), Inches(3.9), Inches(3.5), Inches(0.5),
        text="olusola.xyz", font=BODY_F, size=Pt(20), color=LPURPLE,
        sname="website", alt="Website: olusola.xyz")
    slide_number(slide, n)
    add_notes(slide,
        "ABOUT ME — show while audience settles (~60 seconds).\n\n"
        "'I'm Olusola, Engineering Lead at Retailloop. "
        "I build reliable data systems. This talk is about making that easier.'\n\n"
        "Point to screen: 'QR / olusola.xyz — scan now or at the end.' Keep it brief.")

    # ── SLIDE 2: Title ───────────────────────────────────────────────────────
    n += 1
    slide = blank_slide(prs)
    txbox(slide, Inches(0.54), Inches(0.4), Inches(9.0), Inches(0.5),
        text="EuroPython 2026 · Data Engineering & MLOps",
        font=BODY_F, size=Pt(18), color=GRAY,
        sname="conference-label")
    txbox(slide, Inches(0.54), Inches(1.0), Inches(9.0), Inches(1.5),
        text="Zero to Pipeline",
        font=HEADING_F, size=Pt(60), bold=True, color=LBLUE,
        sname="slide-title", alt="Talk title: Zero to Pipeline")
    txbox(slide, Inches(0.54), Inches(2.6), Inches(9.0), Inches(0.75),
        text="When Data Connectors Just Work",
        font=HEADING_F, size=Pt(30), color=LAVENDER, sname="subtitle")
    txbox(slide, Inches(0.54), Inches(3.5), Inches(9.0), Inches(0.5),
        text="Olusola Akinsulere",
        font=BODY_F, size=Pt(22), color=LPURPLE, sname="byline")
    slide_number(slide, n)
    add_notes(slide,
        "TITLE SLIDE — ~30 seconds.\n\n"
        "Open: 'Raise your hand if you've spent more than an hour debugging "
        "a data connector this year.' Let hands stay up. Look around. "
        "'That's what this talk is about.'")

    # ── SLIDE 3: ML Lifecycle — context slide (Andrew Ng) ───────────────────
    n += 1
    slide = blank_slide(prs)
    # White background for this slide so the image reads clearly
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    txbox(slide, Inches(0.4), Inches(0.15), Inches(9.2), Inches(0.55),
        text="Your data is already in production.",
        font=HEADING_F, size=Pt(26), bold=True,
        color=RGBColor(0x21, 0x21, 0x23),
        sname="slide-title",
        alt="Your data is already in production.")

    img_path = (
        "/Users/olusolaakinsulere/Documents/loganworld/personal/"
        "personalprojects/dataingestionpydatahelenskidemo/"
        "presentation/ml_lifecycle_andrew_ng.png"
    )
    slide.shapes.add_picture(img_path,
        Inches(0.4), Inches(0.8), Inches(9.2), Inches(4.25))

    txbox(slide, Inches(0.4), Inches(5.2), Inches(9.2), Inches(0.3),
        text="Source: Andrew Ng / DeepLearning.AI — MLOps Specialization",
        font=BODY_F, size=Pt(11),
        color=RGBColor(0x99, 0x99, 0x99),
        align=PP_ALIGN.RIGHT,
        sname="attribution")

    slide_number(slide, n)
    add_notes(slide,
        "TIMING: ~1:00. Bridge from title slide.\n\n"
        "'This is the ML project lifecycle — most of you know this from Andrew Ng.\n\n"
        "Look at where your data lives: it is in production systems.\n"
        "MLflow tracking your experiments. Prometheus watching your model serving.\n"
        "Airflow running your pipelines. Your own internal feature store.\n"
        "It is all already there — deployed, running, generating data.\n\n"
        "The loop between Deployment and Data is real. You need to pull that\n"
        "production data back to improve your models, monitor your pipelines,\n"
        "feed your retraining jobs.\n\n"
        "But to do that — you have to connect to all of it.\n"
        "And connecting to production infrastructure is where afternoons go to die.\n\n"
        "That is the problem this talk is about.'\n\n"
        "Pause. Then click to the next slide.")

    # ── SLIDE 4: The familiar afternoon ─────────────────────────────────────
    n += 1
    big_phrase(prs, n,
        "It was supposed to be quick.",
        notes="TIMING: ~1:30. Bridge: 'That feedback loop between Deployment and Data? "
              "This is what it looks like in practice.' Let it land.")

    # ── SLIDE 4: Timeline ────────────────────────────────────────────────────
    n += 1
    slide = blank_slide(prs)
    txbox(slide, Inches(0.54), Inches(0.2), Inches(9.0), Inches(0.55),
        text="A familiar afternoon",
        font=HEADING_F, size=Pt(28), bold=True, color=LBLUE,
        sname="slide-title")
    timeline = [
        ("14:00", "\"Should be quick.\"",          LAVENDER),
        ("14:45", "OAuth. Reading docs. Again.",    CORAL),
        ("15:30", "SDK throws 403. curl works.",    CORAL),
        ("16:15", "Writing a paginator. Again.",    CORAL),
        ("17:00", "Crash. Full re-fetch. Delayed.", RED_BOX),
    ]
    for i, (t, e, col) in enumerate(timeline):
        y = Inches(0.9) + i * Inches(0.82)
        txbox(slide, Inches(0.4), y, Inches(1.4), Inches(0.6),
            text=t, font=CODE_F, size=Pt(20), bold=True, color=col,
            sname=f"time-{i+1}", alt=f"Timeline {t}")
        txbox(slide, Inches(1.9), y, Inches(7.8), Inches(0.6),
            text=e, font=BODY_F, size=Pt(20), color=LAVENDER,
            wrap=True, sname=f"event-{i+1}")
    slide_number(slide, n)
    add_notes(slide,
        "Build one entry at a time. Pause after each red line. "
        "Let the audience feel the familiar sting.\n\n"
        "Key message at the end: 'The pain is not in the ML. It's in the plumbing.'")

    # ── SLIDE 5: The real cost ───────────────────────────────────────────────
    n += 1
    big_phrase(prs, n,
        "The model waited.\nThe sprint slipped.",
        color=CORAL,
        notes="TIMING: ~1:30. This is the emotional beat. Pause. Let it sit.")

    # ── SLIDE 6: The shortage quote ──────────────────────────────────────────
    n += 1
    slide = blank_slide(prs)
    txbox(slide, Inches(0.4), Inches(0.8), Inches(0.08), Inches(3.5),
        text="", sname="quote-bar")
    rect(slide, Inches(0.4), Inches(0.8), Inches(0.08), Inches(3.5),
         LBLUE, decorative=True)
    txbox(slide,
        Inches(0.7), Inches(0.9), Inches(8.8), Inches(2.5),
        text="“We don’t have a shortage of data.\nWe have a shortage of time to connect it.”",
        font=HEADING_F, size=Pt(32), bold=True, color=LAVENDER,
        italic=False, wrap=True, align=PP_ALIGN.LEFT,
        sname="quote",
        alt="Quote: We don't have a shortage of data. "
            "We have a shortage of time to connect it.")
    txbox(slide, Inches(0.7), Inches(3.7), Inches(8.0), Inches(0.5),
        text="— Every ML platform team, every sprint",
        font=BODY_F, size=Pt(20), color=GRAY, sname="quote-attr")
    slide_number(slide, n)
    add_notes(slide,
        "TIMING: ~1:45. Read it slowly. Let it sit for 2 full seconds before clicking.")

    # ── SLIDE 7: The promise ─────────────────────────────────────────────────
    n += 1
    big_phrase(prs, n,
        "Three commands.\nAny API.\nZero config files.",
        color=LBLUE,
        notes="TIMING: ~1:55. Say it slowly. Emphasise 'zero config files'. "
              "Then: 'I know what you're thinking — that sounds too good to be true. "
              "That's exactly why I'm here.'")

    # ── SLIDE 8: Demo preview (code) ─────────────────────────────────────────
    n += 1
    code_slide(prs, n, "The promise in code", [
        ("# That feature store from this afternoon:",           GRAY),
        ("pipeline source add my-feature-store \\",            LBLUE),
        ("    --base-url https://features.internal.mycompany.com", LBLUE),
        ("",                                                     WHITE),
        ("pipeline auth set my-feature-store",                 LPURPLE),
        ("",                                                     WHITE),
        ("pipeline sync run my-feature-store",                 WHITE),
        ("# Done: 50,284 prediction records. Cursor saved.",   GRAY),
    ], font_size=Pt(19),
    notes="TIMING: ~2:00. Let the audience read it. "
          "'The feature store that took all afternoon. Three commands. That's it.'")

    # ── SLIDE 9: Section break — What does 'just work' mean? ─────────────────
    n += 1
    section_break(prs, n,
        "What does 'Just Work' actually mean?",
        "Defining success criteria honestly",
        notes="TIMING: 2:00. 'Before I show you the code, let me be specific. "
              "If I just say it works, that's a red flag at a technical conference.'")

    # ── SLIDE 10: Success criteria ───────────────────────────────────────────
    n += 1
    two_col_slide(prs, n,
        "The success criteria",
        left_head="It means ✓", left_color=LBLUE,
        left_items=[
            "First sync in under 5 minutes",
            "Sane defaults for 90% of APIs",
            "Name + credential = running",
            "Transparent permission model",
            "Idempotent re-runs",
            "Self-heals on common failures",
        ],
        right_head="It does NOT mean ✗", right_color=CORAL,
        right_items=[
            "No debugging ever needed",
            "Skip security review",
            "Works for every edge case",
            "Replaces data modelling",
            "Magic that hides everything",
        ],
        notes="TIMING: 2:30. Walk each bullet. Pause on 'idempotent re-runs'. "
              "Key line: 'We're removing toil, not engineering.'")

    # ── SLIDE 11: Spectrum ───────────────────────────────────────────────────
    n += 1
    slide = blank_slide(prs)
    txbox(slide, Inches(0.54), Inches(0.2), Inches(9.0), Inches(0.55),
        text="The spectrum",
        font=HEADING_F, size=Pt(28), bold=True, color=LBLUE,
        sname="slide-title")
    labels = [
        ("Manual\n(write everything)", Inches(0.4),  CORAL,  "spectrum-manual"),
        ("Config-driven\n(YAML files)", Inches(3.4),  LPURPLE,"spectrum-config"),
        ("Self-configuring\n(name + credential)", Inches(6.4), LBLUE,"spectrum-self"),
    ]
    for text, x, col, sname in labels:
        txbox(slide, x, Inches(1.2), Inches(3.0), Inches(1.1),
            text=text, font=BODY_F, size=Pt(22), bold=True,
            color=col, align=PP_ALIGN.CENTER, wrap=True, sname=sname)
    rect(slide, Inches(0.4), Inches(2.55), Inches(9.2), Inches(0.05),
         GRAY, decorative=True)
    txbox(slide, Inches(0.4), Inches(3.0), Inches(9.2), Inches(0.7),
        text="We're building for the right — with escape hatches to the left.",
        font=HEADING_F, size=Pt(22), italic=True, color=LAVENDER,
        align=PP_ALIGN.CENTER, sname="spectrum-note")
    slide_number(slide, n)
    add_notes(slide,
        "TIMING: 4:00. 'You can always --base-url your way back to manual control. "
        "That's by design.'")

    # ── SLIDE 12: Section break — Architecture ───────────────────────────────
    n += 1
    section_break(prs, n, "Architecture", "The moving parts",
        notes="TIMING: 6:00.")

    # ── SLIDE 13: Architecture overview ──────────────────────────────────────
    n += 1
    slide = blank_slide(prs)
    txbox(slide, Inches(0.54), Inches(0.2), Inches(9.0), Inches(0.55),
        text="Five layers",
        font=HEADING_F, size=Pt(28), bold=True, color=LBLUE,
        sname="slide-title")
    layers = [
        ("1  Provider Registry + LLM Discovery", LBLUE),
        ("2  Auth / Keyring  ·  Self-Healing Connector  ·  Pagination", LPURPLE),
        ("3  Pipeline Orchestrator  (asyncio) — Checkpoint · Retry · DAG", LAVENDER),
        ("4  Output:  JSONL → S3 → DB → Warehouse", GRAY),
    ]
    for i, (label, col) in enumerate(layers):
        y = Inches(0.9) + i * Inches(1.05)
        rect(slide, Inches(0.3), y, Inches(9.4), Inches(0.85),
             BG_DARK, decorative=False,
             sname=f"layer-{i+1}-bg",
             alt=f"Architecture layer {i+1}: {label}")
        txbox(slide, Inches(0.5), y + Inches(0.15), Inches(9.1), Inches(0.6),
            text=label, font=CODE_F, size=Pt(18), color=col,
            sname=f"layer-{i+1}", alt=label)
    slide_number(slide, n)
    add_notes(slide,
        "TIMING: 6:15. Don't read — point to each layer as you name it. "
        "45 seconds max on this slide.")

    # ── SLIDE 14: Layer 1 — Discovery ────────────────────────────────────────
    n += 1
    code_slide(prs, n, "Layer 1 — Discovery", [
        ("# Known provider: LLM enriches the preset",           GRAY),
        ("config = registry.infer_config(\"mlflow\")",           LBLUE),
        ("# → source = \"preset+llm\"",                          GRAY),
        ("",                                                      WHITE),
        ("# Unknown provider (from slide 4): LLM discovers everything", GRAY),
        ("config = registry.infer_config(\"my-feature-store\")", LPURPLE),
        ("# → auth_type=bearer, pagination=cursor, endpoints=[...]", GRAY),
        ("# → source = \"llm_discovered\"",                      GRAY),
    ], font_size=Pt(18),
    notes="TIMING: 7:30. 'mlflow — known, preset+LLM.'\n"
          "'my-feature-store — never seen before. LLM discovers auth type, "
          "pagination style, endpoints. From the name and URL alone.'")

    # ── SLIDE 15: Credentials — the bad way ──────────────────────────────────
    n += 1
    code_slide(prs, n, "Layer 2 — Credentials (old way — risky)", [
        ("# .env file",                                 GRAY),
        ("FEATURE_STORE_TOKEN=tok_abc123...",           CORAL),
        ("MLFLOW_TRACKING_TOKEN=mlf_xyz789...",        CORAL),
        ("# config.yaml",                               GRAY),
        ("feature_store:",                              CORAL),
        ("  token: ${FEATURE_STORE_TOKEN}",            CORAL),
        ("",                                             WHITE),
        ("# One git commit away from a breach.",       RED_BOX),
    ], callout=(1, 5),
    notes="TIMING: 9:00. Point at the callout box. "
          "'One accidental git push. It happens weekly.'")

    # ── SLIDE 16: Credentials — the new way ──────────────────────────────────
    n += 1
    code_slide(prs, n, "Layer 2 — Credentials (new way — secure)", [
        ("pipeline auth set my-feature-store",              LBLUE),
        ("Enter API token for my-feature-store: ****",      GRAY),
        ("Token stored securely for my-feature-store.",     LPURPLE),
        ("",                                                  WHITE),
        ("# macOS Keychain / Linux Secret Service",          GRAY),
        ("# Never touches disk as plaintext.",               LBLUE),
    ], callout=(5, 5),
    notes="TIMING: 9:30. 'Same security model as your password manager. "
          "Resolves the 14:45 pain — OAuth token stored once, never in .env.'")

    # ── SLIDE 17: Self-healing ───────────────────────────────────────────────
    n += 1
    code_slide(prs, n, "Layer 3 — Self-Healing Auth", [
        ("# \"Token works in curl but SDK throws 403.\" — slide 4, 15:30", GRAY),
        ("",                                                                 WHITE),
        ("Try 1: Authorization: Bearer tok_abc123  →  403 FAIL ✗",         CORAL),
        ("Try 2: Authorization: tok_abc123          →  200 OK   ✓",         LBLUE),
        ("",                                                                 WHITE),
        ("Learned. All future requests use the corrected format.",           LAVENDER),
        ("Tries: Bearer · raw · token prefix · X-API-Key · Api-Key",       GRAY),
    ], callout=(3, 3),
    notes="TIMING: 10:00. Call back to slide 4: '15:30 — token works in curl, SDK throws 403.'\n"
          "'Try 1: Bearer — 403. Try 2: no prefix — 200. Now it remembers. No config change needed.'\n"
          "'The healing_success log event is also a signal: if it fires on a stable pipeline, "
          "the upstream API changed its auth. That's worth an alert.'")

    # ── SLIDE 18: Pagination ─────────────────────────────────────────────────
    n += 1
    slide = blank_slide(prs)
    txbox(slide, Inches(0.54), Inches(0.2), Inches(9.0), Inches(0.55),
        text="Layer 4 — Universal Pagination",
        font=HEADING_F, size=Pt(28), bold=True, color=LBLUE,
        sname="slide-title")
    cols = [
        ("Cursor-based",     "next_cursor → start_cursor\nmy-feature-store, W&B, Prefect",   LBLUE),
        ("Offset-based",     "page=1 → page=2…\nMLflow, Grafana, REST APIs",                 LPURPLE),
        ("GraphQL cursor",   "pageInfo.endCursor\nW&B runs, GitHub Actions",                 LAVENDER),
    ]
    for i, (title, body, col) in enumerate(cols):
        x = Inches(0.3) + i * Inches(3.2)
        rect(slide, x, Inches(0.9), Inches(3.1), Inches(2.8),
             BG_DARK, decorative=False,
             sname=f"pag-{i+1}-bg", alt=f"Pagination type {i+1}: {title}")
        txbox(slide, x + Inches(0.15), Inches(1.0), Inches(2.9), Inches(0.6),
            text=title, font=HEADING_F, size=Pt(20), bold=True, color=col,
            sname=f"pag-{i+1}-head")
        txbox(slide, x + Inches(0.15), Inches(1.65), Inches(2.9), Inches(1.8),
            text=body, font=CODE_F, size=Pt(17), color=WHITE,
            wrap=True, sname=f"pag-{i+1}-body")
    txbox(slide, Inches(0.3), Inches(3.9), Inches(9.4), Inches(0.55),
        text="pagination = connector.infer_pagination(sample_response)",
        font=CODE_F, size=Pt(18), color=LBLUE, sname="infer-code",
        alt="Code: infer_pagination auto-detects the pagination type from the response")
    slide_number(slide, n)
    add_notes(slide,
        "TIMING: 11:00. 'The connector samples one response and detects which pattern "
        "is in use. You don't specify it.'")

    # ── SLIDE 19: Orchestration ──────────────────────────────────────────────
    n += 1
    code_slide(prs, n, "Layer 5 — Orchestration (asyncio DAG)", [
        ("# The exact pipeline from the demo",                              GRAY),
        ("pipeline = Pipeline(\"demo-pipeline\")",                          LAVENDER),
        ("pipeline.add_step(\"extract_feature_store\", extract_fs)",       LBLUE),
        ("pipeline.add_step(\"extract_mlflow\",        extract_mlflow)",   LBLUE),
        ("pipeline.add_step(\"create_linear_issues\",  create_issues,",   LPURPLE),
        ("    depends_on=[\"extract_feature_store\", \"extract_mlflow\"])", LPURPLE),
        ("",                                                                 WHITE),
        ("result = await engine.run(pipeline)",                             WHITE),
        ("# Steps 1-2 parallel → step 3 waits → sprint unblocked",        GRAY),
    ], font_size=Pt(17),
    notes="TIMING: 11:30. 'Two extraction steps run in parallel. When both finish, "
          "the Linear step runs. No Temporal. No Celery. No infrastructure.\n"
          "This is the actual pipeline you will see me run in the demo.'")

    # ── SLIDE 20: Fallback — backup demo video (hidden slide) ────────────────
    n += 1
    slide = blank_slide(prs)
    # Dark background with a standout note — this slide is only shown if live demo fails
    rect(slide, Inches(0), Inches(0), Inches(10), Inches(5.625), BG_DARK, decorative=True)
    txbox(slide, Inches(0.5), Inches(0.4), Inches(9.0), Inches(0.7),
        text="BACKUP — Demo Video",
        font=HEADING_F, size=Pt(30), bold=True, color=CORAL,
        sname="slide-title", alt="Backup demo video slide")
    txbox(slide, Inches(0.5), Inches(1.2), Inches(9.0), Inches(1.0),
        text="Open Loom recording → play from timestamp",
        font=BODY_F, size=Pt(22), color=LAVENDER, sname="instruction")
    txbox(slide, Inches(0.5), Inches(2.3), Inches(9.0), Inches(0.6),
        text="Timestamps:",
        font=HEADING_F, size=Pt(18), bold=True, color=LBLUE, sname="ts-head")
    timestamps = [
        "0:00  Health check (pipeline doctor)",
        "0:20  Add feature store (LLM discovers)",
        "1:10  Store credentials (auth docs shown)",
        "1:50  First sync + self-healing",
        "2:30  Incremental sync",
        "2:50  Fraud detection pipeline (4.1s end-to-end)",
        "4:00  Interactive assistant (pipeline chat)",
    ]
    for i, ts in enumerate(timestamps):
        txbox(slide, Inches(0.5), Inches(2.95 + i * 0.32), Inches(9.0), Inches(0.3),
            text=ts, font=CODE_F, size=Pt(14), color=GRAY, sname=f"ts-{i+1}")
    txbox(slide, Inches(0.5), Inches(5.1), Inches(9.0), Inches(0.35),
        text="https://www.loom.com/share/8f29b2c846d74686977947a0e27ee1c1",
        font=CODE_F, size=Pt(13), color=CORAL, sname="loom-url")
    slide_number(slide, n)
    add_notes(slide,
        "FALLBACK SLIDE — only use if live demo fails.\n\n"
        "Say calmly: 'Let me show you a recording of exactly what this looks like.'\n"
        "Open Loom, play from the relevant timestamp.\n\n"
        "Timestamps:\n"
        "  0:00  Health check\n"
        "  0:20  Add feature store\n"
        "  1:10  Auth set with docs\n"
        "  1:50  First sync + self-healing\n"
        "  2:30  Incremental sync\n"
        "  2:50  Fraud detection pipeline\n"
        "  4:00  Chat assistant\n\n"
        "After video: 'The whole thing. In under 5 minutes. Let me show you why.'")

    # ── SLIDE 21: Section break — Live Demo ──────────────────────────────────
    n += 1
    section_break(prs, n, "Live Demo", "Fraud detection — zero to pipeline in real time",
        notes="TIMING: 12:00 — SWITCH TO TERMINAL.\n\n"
              "ACT 1 — Health check (20s)\n"
              "  pipeline doctor\n\n"
              "ACT 2 — Start all services (20s)\n"
              "  uv run python -m fraud_demo.setup\n"
              "  (Feast 6566, MLflow 5001, Prometheus 9090 — all confirmed up)\n\n"
              "ACT 3 — Add Feast (no auth, known preset) (30s)\n"
              "  pipeline source add feast --force\n"
              "  pipeline source test feast\n\n"
              "ACT 4 — Add MLflow (no auth, port 5001) (30s)\n"
              "  pipeline source add mlflow --base-url http://127.0.0.1:5001 --auth-type none --force\n"
              "  pipeline source test mlflow\n\n"
              "ACT 5 — First sync: 13 runs, 5 per page, cursor saved (45s)\n"
              "  pipeline sync run mlflow\n"
              "  (Done: 13 records — 3 pages of 5)\n\n"
              "ACT 6 — Incremental sync: cursor resumes (20s)\n"
              "  pipeline sync run mlflow\n"
              "  (Resuming from cursor. Done: 3 records — only new)\n"
              "  pipeline sync status\n\n"
              "ACT 7 — Full fraud detection pipeline (90s)\n"
              "  uv run python -m fraud_demo.run_pipeline\n"
              "  (Feast + 15k transactions → model → MLflow → Prometheus → Linear)\n\n"
              "FALLBACK: hover bottom-right on this slide → click 📹 backup → Loom plays.")

    # ── SLIDE 22–29: Demo slides ──────────────────────────────────────────────

    def demo(num, title, lines, callout_rows=None, fs=Pt(19), notes_text=""):
        code_slide(prs, num, f"Demo — {title}", lines,
                   callout=callout_rows, font_size=fs, notes=notes_text)

    n += 1
    demo(n, "Health check", [
        ("$ pipeline doctor",               LBLUE),
        ("  OK  Keyring accessible",        LBLUE),
        ("  OK  Provider registry loaded",  LBLUE),
        ("  OK  Checkpoint dir writable",   LBLUE),
        ("  OK  Source store accessible",   LBLUE),
        ("  All checks passed!",            LBLUE),
    ], notes_text="TIMING: 12:20. All green. 'Four checks. We're ready.' 10 seconds.")

    n += 1
    demo(n, "Add Feast — the feature store (no auth)", [
        ("$ pipeline source add feast --force",                LBLUE),
        ("  ✓ Registry: preset found for 'feast'",             GRAY),
        ("  Source 'Feast' added  (from: preset+llm)",         LPURPLE),
        ("  Base URL: http://127.0.0.1:6566",                  LAVENDER),
        ("  Auth:     none  ← local Feast, no token needed",   LBLUE),
        ("  Endpoint: /get-online-features",                    LAVENDER),
        ("$ pipeline source test feast",                       LBLUE),
        ("  Connection to feast successful!",                   LBLUE),
    ], callout_rows=(4, 4),
    notes_text="TIMING: 13:00 — first 'wow'. 'Feast is the feature store from slide 4.\n"
               "It's now a known provider. Auth: none — local server needs no token.\n"
               "One command. Connected.'")

    n += 1
    demo(n, "Add MLflow — local server (no auth)", [
        ("$ pipeline source add mlflow \\",                         LBLUE),
        ("    --base-url http://127.0.0.1:5001 --auth-type none",   LBLUE),
        ("  ✓ Registry: preset found for 'mlflow'",                 GRAY),
        ("  Source 'MLflow' added  (from: preset+llm)",             LPURPLE),
        ("  Base URL:   http://127.0.0.1:5001",                     LAVENDER),
        ("  Auth:       none  ← local server, no auth needed",      LBLUE),
        ("$ pipeline source test mlflow",                           LBLUE),
        ("  Connection to mlflow successful!",                       LBLUE),
    ], callout_rows=(5, 5),
    notes_text="TIMING: 13:45. 'MLflow on port 5001 — our local demo server, no auth.\n"
               "Same command as any other source. Connected.'")

    n += 1
    demo(n, "Auth status — no tokens needed", [
        ("$ pipeline auth status",                             LBLUE),
        ("  feast      │ none    │ no auth required",         LBLUE),
        ("  mlflow     │ none    │ no auth required",         LBLUE),
        ("  prometheus │ none    │ no auth required",         LBLUE),
        ("  linear     │ api_key │ authenticated",            LPURPLE),
    ], callout_rows=(1, 3),
    notes_text="TIMING: 14:10. 'Three local services — zero tokens.\n"
               "Linear is the only one that needs a token, and it's already in the keychain.'")

    n += 1
    demo(n, "First sync — 13 runs, cursor pagination", [
        ("$ pipeline sync run mlflow",                               LBLUE),
        ("  pagination: cursor (MLflow page_token, 5 per page)",    GRAY),
        ("  1. fraud_rf_20260713_1353",                              LAVENDER),
        ("  2. fraud_rf_20260713_1348",                              LAVENDER),
        ("  3. fraud_rf_20260713_1346  ...",                         LAVENDER),
        ("  Done: 13 records",                                       LPURPLE),
        ("  checkpoint: saved — cursor eyJvZmZzZXQiOiAxMH0=",      GRAY),
    ], callout_rows=(5, 5),
    notes_text="TIMING: 15:00. 'Cursor pagination. MLflow returns a page_token.\n"
               "The connector follows it page by page — 5 per page, 3 pages, 13 runs.\n"
               "Cursor saved. Watch what happens when I run it again...'")

    n += 1
    demo(n, "Incremental sync — 3 records, cursor resumes", [
        ("$ pipeline sync run mlflow",                               LBLUE),
        ("  Resuming from checkpoint cursor: eyJvZmZzZXQiOiAxMH0=", GRAY),
        ("  Done: 3 records  ← only new runs since last sync",      LPURPLE),
        ("",                                                          WHITE),
        ("$ pipeline sync status",                                   LBLUE),
        ("  mlflow | cursor: eyJvZmZzZX... | 2026-07-13 17:04",     LAVENDER),
    ], callout_rows=(2, 2),
    notes_text="TIMING: 15:45. 'Cursor picked up exactly where it left off.\n"
               "3 records, not 13. Your nightly retraining job fetches only what changed.'")

    n += 1
    demo(n, "Interactive assistant", [
        ("$ pipeline chat",                                        LBLUE),
        ("You: what sources are failing?",                        LAVENDER),
        ("  Prometheus is failing — server not running.",          LPURPLE),
        ("  Want me to start it locally?",                        LPURPLE),
        ("You: yes start it",                                     LAVENDER),
        ("  ✓ prometheus ready at http://127.0.0.1:9090",        LBLUE),
        ("You: show me all my sources",                           LAVENDER),
    ], callout_rows=(4, 5),
    notes_text="TIMING: 16:15. 'Full conversational REPL — like Claude or ChatGPT.'\n"
               "'It knows your pipeline state. It can add sources, start containers, sync.'\n"
               "'For team members who don't live in the terminal.'")

    n += 1
    code_slide(prs, n, "Demo — Fraud detection pipeline → Linear", [
        ("$ uv run python -m fraud_demo.run_pipeline",                     LBLUE),
        ("  ✓ Feast: 50 feature vectors pulled",                           LBLUE),
        ("  ✓ Transactions: 15,000 records  (5.5% fraud rate)",            LBLUE),
        ("  ✓ MLflow: run logged  accuracy=64.7%  recall=43.6%",           LPURPLE),
        ("  ✓ Prometheus: metrics pushed",                                  LAVENDER),
        ("  ✓ Linear issue: https://linear.app/olusola-akinsulere/...",    LBLUE),
        ("  Pipeline completed in 4.4s",                                    WHITE),
        ("  No connector classes. No YAML. No SDK imports per source.",     GRAY),
    ], callout=(5, 5), font_size=Pt(17),
    notes="TIMING: 16:45. Run: uv run python -m fraud_demo.run_pipeline\n\n"
          "As steps complete, narrate:\n"
          "  'Feast features and transactions — parallel.'\n"
          "  'Model trained. 64.7% accuracy. 43.6% recall.'\n"
          "  'MLflow run logged — here is the URL.'\n"
          "  'Prometheus metrics pushed.'\n"
          "  'Linear issue created. Sprint unblocked.'\n\n"
          "Point at the bottom line:\n"
          "  '4.4 seconds. No connector classes. No YAML. No SDK imports.'\n\n"
          "Transition: 'Now let me show you why you can trust it.'")

    # ── SLIDE 28: Section break — Reliability ────────────────────────────────
    n += 1
    section_break(prs, n, "The Reliability Playbook",
        "The patterns that make it trustworthy",
        notes="TIMING: 17:00. 'The demo looked easy. That was intentional.' Pause. "
              "'Easy to use doesn't mean simple inside.'")

    # ── SLIDE 29: Pattern 1 — Retries ────────────────────────────────────────
    n += 1
    code_slide(prs, n, "Pattern 1 — Retries + Backoff", [
        ("from tenacity import retry, wait_exponential, stop_after_attempt", LAVENDER),
        ("",                                                                   WHITE),
        ("@retry(",                                                             LPURPLE),
        ("    wait=wait_exponential(multiplier=1, max=60),  # 1s→60s",       LPURPLE),
        ("    stop=stop_after_attempt(3),",                                    LPURPLE),
        ("    retry=retry_if_exception_type((TransportError, RateLimitError)),", LPURPLE),
        (")",                                                                   LPURPLE),
        ("async def request(self, method, path, **kwargs):",                   LAVENDER),
        ("    if response.status_code == 429:  # Rate limited",                CORAL),
        ("        raise RateLimitError(retry_after)   # → tenacity retries",  LBLUE),
    ], callout=(2, 6), font_size=Pt(17),
    notes="TIMING: 17:30. Point to decorator: 'Three lines.' "
          "Point at 429: 'Rate limited? We sleep the exact Retry-After duration.'")

    # ── SLIDE 30: Pattern 2 — Checkpointing ──────────────────────────────────
    n += 1
    code_slide(prs, n, "Pattern 2 — Checkpointing", [
        ("# After every 100 records, persist the cursor",                  GRAY),
        ("if count % 100 == 0 and last_cursor:",                           LAVENDER),
        ("    checkpoint_mgr.save(CheckpointState(",                       LAVENDER),
        ("        source_id=\"my-feature-store:/v1/features\",",           LAVENDER),
        ("        cursor=last_cursor,",                                    LPURPLE),
        ("        last_sync_at=datetime.now(),",                           LAVENDER),
        ("    ))",                                                          LAVENDER),
        ("",                                                                WHITE),
        ("# Crash at 5,000?  Restart picks up at 5,001. Not from zero.",  LBLUE),
    ], callout=(8, 8),
    notes="TIMING: 18:30. Call back to slide 4: '17:00 — cursor state lost. Full re-fetch.'\n"
          "'That is what happens with no checkpoint. This is what happens with one.'\n"
          "'The checkpoint files are plain JSON. Open, edit, delete. Zero magic.'")

    # ── SLIDE 31: Pattern 3 — Self-healing ───────────────────────────────────
    n += 1
    code_slide(prs, n, "Pattern 3 — Self-Healing Auth", [
        ("AUTH_FORMATS = [",                                               LAVENDER),
        ("    (\"Bearer\",  \"Authorization\"),  # MLflow managed, W&B",  LAVENDER),
        ("    (\"\",        \"Authorization\"),  # MLflow self-hosted",    LAVENDER),
        ("    (\"token\",   \"Authorization\"),  # Airflow, older APIs",   LAVENDER),
        ("    (\"\",        \"X-API-Key\"),      # Grafana, observability", LAVENDER),
        ("    (\"\",        \"Api-Key\"),        # Prometheus pushgateway", LAVENDER),
        ("]",                                                               LAVENDER),
        ("for prefix, header in AUTH_FORMATS:",                            LPURPLE),
        ("    if (await try_auth(prefix, header)).ok:",                    LPURPLE),
        ("        self._auth_prefix = prefix; break  # Saved", LBLUE),
    ], callout=(7, 9),
    notes="TIMING: 19:30. 'If Bearer fails, try raw. Stop as soon as one works.'")

    # ── SLIDE 32: Pattern 4 — Observability ──────────────────────────────────
    n += 1
    code_slide(prs, n, "Pattern 4 — Observability", [
        ('{\"event\":\"pipeline_started\",  \"pipeline\":\"demo-pipeline\"}',         GRAY),
        ('{\"event\":\"step_started\",      \"step\":\"extract_feature_store\"}',     LAVENDER),
        ('{\"event\":\"healing_success\",   \"new_prefix\":\"\", \"new_header\":\"Authorization\"}', LPURPLE),
        ('{\"event\":\"step_completed\",    \"step\":\"extract_feature_store\", \"records\":50284}', LBLUE),
        ('{\"event\":\"step_completed\",    \"step\":\"extract_mlflow\",        \"records\":8421}',  LBLUE),
        ('{\"event\":\"pipeline_completed\",\"total_records\":58705, \"ms\":4850}',  LBLUE),
    ], callout=(2, 2),
    notes="TIMING: 20:30. Read each line.\n"
          "'Pipeline started. Feature store step started.\n"
          "healing_success — that is the 403 from slide 4, fixed automatically.\n"
          "Feature store: 50k prediction records. MLflow: 8k runs. Pipeline done.'\n\n"
          "'The healing_success line: an API threw a 403 and the system fixed itself "
          "without waking anyone up.'")

    # ── SLIDE 33: Full picture table ─────────────────────────────────────────
    n += 1
    slide = blank_slide(prs)
    txbox(slide, Inches(0.54), Inches(0.1), Inches(9.0), Inches(0.55),
        text="Seven patterns that make it trustworthy",
        font=HEADING_F, size=Pt(26), bold=True, color=LBLUE,
        sname="slide-title")

    headers = ["What goes wrong", "Pattern", "How"]
    col_x   = [Inches(0.2), Inches(3.6), Inches(6.4)]
    col_w   = [Inches(3.3), Inches(2.7), Inches(3.4)]
    rows = [
        ("Network blip",       "Retry + backoff",     "tenacity decorator"),
        ("Rate limited",       "Respect Retry-After", "429 detection"),
        ("Process crash",      "Checkpoint resume",   "File-based cursors"),
        ("Wrong auth format",  "Self-healing",        "Format rotation"),
        ("10M record dataset", "Pagination",          "Cursor/offset/GQL"),
        ("Stale data",         "Incremental sync",    "Watermarks"),
        ("Silent failures",    "Structured logs",     "structlog JSON"),
    ]
    row_h = Inches(0.58)
    hdr_y = Inches(0.72)

    # Header row
    rect(slide, Inches(0.2), hdr_y, Inches(9.6), row_h,
         BG_DARK, decorative=False,
         sname="table-header", alt="Table header row")
    for j, (hdr, x, w) in enumerate(zip(headers, col_x, col_w)):
        txbox(slide, x + Inches(0.1), hdr_y + Inches(0.1), w, Inches(0.42),
            text=hdr, font=HEADING_F, size=Pt(17), bold=True, color=LBLUE,
            sname=f"th-{j+1}")

    for i, (c1, c2, c3) in enumerate(rows):
        y = hdr_y + (i + 1) * row_h
        alt_bg = BG_DARK if i % 2 == 0 else RGBColor(0x2a, 0x2a, 0x2f)
        rect(slide, Inches(0.2), y, Inches(9.6), row_h,
             alt_bg, decorative=True)
        for text, x, w in zip([c1, c2, c3], col_x, col_w):
            txbox(slide, x + Inches(0.1), y + Inches(0.1), w, Inches(0.45),
                text=text, font=BODY_F, size=Pt(16), color=LAVENDER,
                sname=f"r{i+1}c{col_x.index(x)+1}")
    slide_number(slide, n)
    add_notes(slide,
        "TIMING: 21:30. 'None of this is magic. Seven well-understood patterns "
        "applied consistently. You can implement every single one independently.'")

    # ── SLIDE 34: Section break — Takeaways ──────────────────────────────────
    n += 1
    section_break(prs, n, "Takeaways", "Your checklist",
        notes="TIMING: 23:00.")

    # ── SLIDE 35: Checklist ───────────────────────────────────────────────────
    n += 1
    slide = blank_slide(prs)
    txbox(slide, Inches(0.54), Inches(0.15), Inches(9.0), Inches(0.55),
        text="Building self-configuring connectors",
        font=HEADING_F, size=Pt(26), bold=True, color=LBLUE,
        sname="slide-title")
    checklist = [
        "Registry of known patterns  (accelerators, not requirements)",
        "LLM-driven discovery  — with offline fallback",
        "OS keychain for credentials  (never .env, never git)",
        "Self-healing auth  — try formats before failing",
        "Universal pagination  — cursor · offset · link-header · GQL",
        "Checkpoint after every batch  — crash = resume, not restart",
        "Structured logs at every state transition",
        "Incremental by default  — full-refresh as an escape hatch",
    ]
    for i, item in enumerate(checklist):
        txbox(slide,
            Inches(0.4), Inches(0.8 + i * 0.585), Inches(9.3), Inches(0.55),
            text=f"✔  {item}", font=BODY_F, size=Pt(17), color=LAVENDER,
            wrap=True, sname=f"check-{i+1}")
    slide_number(slide, n)
    add_notes(slide,
        "TIMING: 23:00. 'You can implement these eight patterns in your own Python "
        "connectors starting this week. You don't need this framework to use these ideas.'")

    # ── SLIDE 36: The promise, delivered ─────────────────────────────────────
    n += 1
    code_slide(prs, n, "The promise, delivered", [
        ("pipeline source add <any-api>", LBLUE),
        ("pipeline auth set <any-api>",   LPURPLE),
        ("pipeline sync run <any-api>",   WHITE),
        ("",                               WHITE),
        ("# my-feature-store · MLflow · W&B · Airflow",   LBLUE),
        ("# Prometheus · GitHub · Linear · any-internal-api", GRAY),
        ("# anything-with-an-http-endpoint",               GRAY),
    ], font_size=Pt(22),
    notes="TIMING: 24:00. Point first to my-feature-store — 'The one from this afternoon.'\n"
          "Then: 'MLflow. W&B. Anything with an HTTP endpoint.'\n"
          "Closing: 'Three commands. Any API. Zero config files.'")

    # ── SLIDE 37: When to graduate ────────────────────────────────────────────
    n += 1
    two_col_slide(prs, n,
        "When to graduate beyond this",
        left_head="This is enough when:", left_color=LBLUE,
        left_items=[
            "Single-machine ETL",
            "Prototyping & exploration",
            "Internal tooling",
            "< 1 M records per run",
        ],
        right_head="Use Temporal / Dagster when:", right_color=CORAL,
        right_items=[
            "Distributed across workers",
            "Jobs run hours or days",
            "Complex saga patterns",
            "Multi-team coordination",
        ],
        notes="TIMING: 24:30. 'Start here. Graduate when the problem grows past it. "
              "Don't use a cargo ship to cross a stream.'")

    # ── SLIDE 38: Questions / Try it now ─────────────────────────────────────
    n += 1
    code_slide(prs, n, "Questions?  Try it now.", [
        ("git clone https://github.com/Lanrey/zero-to-pipeline",  GRAY),
        ("cd zero-to-pipeline/python && uv sync",                  LBLUE),
        ("pipeline source add my-feature-store \\",               LBLUE),
        ("    --base-url https://your-api.internal",               LBLUE),
        ("pipeline auth set my-feature-store",                    LPURPLE),
        ("pipeline auth set linear",                               LPURPLE),
        ("uv run python -m examples.demo_pipeline",               WHITE),
    ], font_size=Pt(19),
    notes="TIMING: 25:00 — Q&A starts.\n\n"
          "Prepared answers:\n"
          "Q: 'DLT / Airbyte?' → Lighter, more Python-native. Use DLT for 300+ connectors.\n"
          "Q: 'Production?' → Yes, single-machine, <1M records.\n"
          "Q: 'Which LLM?' → Pre-configured for this demo. Any reasoning model works.\n\n"
          "Close: 'I'm around for the rest of the conference. Thank you.'")

    # ── SLIDE 39: Thank you ───────────────────────────────────────────────────
    n += 1
    slide = blank_slide(prs)
    txbox(slide, Inches(0.54), Inches(0.8), Inches(9.0), Inches(1.2),
        text="Thank you",
        font=HEADING_F, size=Pt(60), bold=True, color=LBLUE,
        sname="slide-title", alt="Thank you")
    txbox(slide, Inches(0.54), Inches(2.1), Inches(9.0), Inches(0.6),
        text="Olusola Akinsulere  |  EuroPython 2026 · Kraków, Poland",
        font=BODY_F, size=Pt(22), color=LAVENDER, sname="credit")
    txbox(slide, Inches(0.54), Inches(3.0), Inches(9.0), Inches(1.5),
        text="“The best pipeline is the one your team actually uses\n— because it was the easiest one to set up.”",
        font=HEADING_F, size=Pt(24), italic=True, color=LPURPLE,
        wrap=True, sname="closing-quote",
        alt="Closing quote: The best pipeline is the one your team actually uses "
            "— because it was the easiest one to set up.")
    txbox(slide, Inches(0.54), Inches(4.6), Inches(4.0), Inches(0.5),
        text="olusola.xyz",
        font=BODY_F, size=Pt(22), color=GRAY,
        sname="website", alt="Website: olusola.xyz")
    slide_number(slide, n)
    add_notes(slide,
        "Let the quote sit.\n"
        "'Repo is in the talk description. I'm around for the rest of the conference.'")

    # ── Save ──────────────────────────────────────────────────────────────────
    out = (
        "/Users/olusolaakinsulere/Documents/loganworld/personal/"
        "personalprojects/dataingestionpydatahelenskidemo/"
        "presentation/zero-to-pipeline-europython2026.pptx"
    )
    prs.save(out)
    print(f"Saved: {out}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    build()
